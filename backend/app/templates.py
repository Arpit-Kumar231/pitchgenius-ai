"""Template ingestion: store uploaded .pptx files and extract STYLE/STRUCTURE only.

We never read content/data from the source decks. We only inspect:
  - slide masters & layout names (structure)
  - theme colors & default fonts (style)
  - slide count
"""
from __future__ import annotations

import logging
import os
import uuid
from typing import Any

from pptx import Presentation

logger = logging.getLogger("pitchbook.templates")

TEMPLATES_DIR = os.environ.get("TEMPLATES_DIR", "/tmp/pitchbook_templates")
os.makedirs(TEMPLATES_DIR, exist_ok=True)

# In-memory index. {id: {id, name, path, slide_count, layouts, fonts, theme_colors}}
TEMPLATES: dict[str, dict[str, Any]] = {}


def _extract_style(path: str) -> dict[str, Any]:
    prs = Presentation(path)
    layouts = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            layouts.append(layout.name)
    fonts: set[str] = set()
    # Sniff fonts from master placeholders (no slide content read)
    for master in prs.slide_masters:
        for shape in master.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font and run.font.name:
                            fonts.add(run.font.name)
    return {
        "slide_count": len(prs.slides),
        "layouts": layouts[:20],
        "fonts": sorted(fonts)[:10],
        "slide_width_in": float(prs.slide_width) / 914400.0,
        "slide_height_in": float(prs.slide_height) / 914400.0,
    }


def save_template(filename: str, data: bytes) -> dict[str, Any]:
    tid = uuid.uuid4().hex[:10]
    safe = os.path.basename(filename) or f"template_{tid}.pptx"
    if not safe.lower().endswith(".pptx"):
        raise ValueError("Only .pptx files are supported")
    path = os.path.join(TEMPLATES_DIR, f"{tid}__{safe}")
    with open(path, "wb") as f:
        f.write(data)
    try:
        style = _extract_style(path)
    except Exception:
        logger.exception("Failed to parse template %s", safe)
        os.remove(path)
        raise
    record = {"id": tid, "name": safe, "path": path, **style}
    TEMPLATES[tid] = record
    logger.info("Ingested template %s (%s) — %d layouts", safe, tid, len(style["layouts"]))
    return {k: v for k, v in record.items() if k != "path"}


def list_templates() -> list[dict[str, Any]]:
    return [{k: v for k, v in t.items() if k != "path"} for t in TEMPLATES.values()]


def get_template_path(tid: str) -> str | None:
    rec = TEMPLATES.get(tid)
    return rec["path"] if rec else None