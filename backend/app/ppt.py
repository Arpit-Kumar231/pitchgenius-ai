"""Pitchbook PPT builder.

Generates an investment-banking style deck. The structure mirrors the
common pitchbook flow:
  1. Cover
  2. Executive summary
  3. Situation overview
  4. Market & industry
  5. Client snapshot (CRM)
  6. Competitor landscape
  7. Financial profile & comps
  8. Strategic alternatives
  9. Why us
 10. Next steps

When you ingest your real templates later, swap `_blank_deck()` for
`Presentation(template_path)` and reuse its slide layouts.
"""
from __future__ import annotations

import logging
import os
import uuid
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

logger = logging.getLogger("pitchbook.ppt")

OUT_DIR = os.environ.get("PPT_OUT_DIR", "/tmp/pitchbooks")
os.makedirs(OUT_DIR, exist_ok=True)

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
GOLD = RGBColor(0xC8, 0xA2, 0x5B)
LIGHT = RGBColor(0xF4, 0xF1, 0xEA)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)


def _blank_deck(template_path: str | None = None) -> Presentation:
    """Return a new Presentation. If a template path is given, inherit its
    slide masters, layouts, theme colors and fonts. We CLEAR any pre-existing
    slides from the template so we never leak source-deck content."""
    if template_path and os.path.isfile(template_path):
        try:
            prs = Presentation(template_path)
            # Drop any slides that came with the template — keep masters/layouts only
            xml_slides = prs.slides._sldIdLst  # type: ignore[attr-defined]
            for sld in list(xml_slides):
                xml_slides.remove(sld)
            logger.info("Using template %s (masters=%d, layouts inherited)",
                        os.path.basename(template_path), len(prs.slide_masters))
            return prs
        except Exception:
            logger.exception("Template load failed, falling back to blank deck")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False


def _text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK, align=None):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tx


def _bar(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _section_header(slide, title: str):
    _bar(slide, 0, 0, 13.333, 0.9, NAVY)
    _bar(slide, 0, 0.9, 13.333, 0.05, GOLD)
    _text(slide, 0.5, 0.18, 12, 0.6, title, size=26, bold=True, color=RGBColor(255, 255, 255))


def _bullets(slide, x, y, w, h, items, *, size=16):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = DARK
        r.font.name = "Calibri"
        p.space_after = Pt(6)


# ---------------------------------------------------------------------------
def build_pitchbook(
    *,
    topic: str,
    client: str,
    research: dict[str, Any],
    crm: dict[str, Any],
    competitors: dict[str, Any],
    financials: dict[str, Any],
    template_path: str | None = None,
) -> str:
    prs = _blank_deck(template_path)
    # Pick a blank-ish layout. python-pptx convention: index 6 is "Blank" in the
    # default template; for custom templates we fall back to the LAST layout
    # (typically the simplest), which is safer than assuming index 6 exists.
    layouts = prs.slide_layouts
    blank = layouts[6] if len(layouts) > 6 else layouts[len(layouts) - 1]

    # 1 Cover ----------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s, NAVY)
    _bar(s, 0.5, 3.2, 1.2, 0.08, GOLD)
    _text(s, 0.5, 3.4, 12, 1.0, topic, size=44, bold=True, color=RGBColor(255, 255, 255))
    _text(s, 0.5, 4.4, 12, 0.6, f"Prepared for {client}", size=22, color=LIGHT)
    _text(s, 0.5, 6.7, 12, 0.4, "Strictly private & confidential", size=12, color=GOLD)

    # 2 Executive summary ----------------------------------------------------
    s = prs.slides.add_slide(blank)
    _section_header(s, "Executive Summary")
    _bullets(
        s,
        0.6,
        1.3,
        12,
        5.5,
        [
            f"{client} is positioned to capitalize on a ${research.get('market_size_usd_bn','—')}B market growing at {research.get('cagr_pct','—')}% CAGR.",
            f"Current relationship: {crm.get('tier','—')} tier, {crm.get('wallet_share_pct','—')}% wallet share.",
            f"Financial profile: {financials.get('ebitda_margin_pct','—')}% EBITDA margin, {financials.get('net_leverage_x','—')}x net leverage.",
            "Recommended next step: explore a strategic financing / advisory mandate (see slide 8).",
        ],
        size=18,
    )

    # 3 Situation ------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _section_header(s, "Situation Overview")
    _text(s, 0.6, 1.3, 12, 0.5, f"Topic: {topic}", size=18, bold=True, color=NAVY)
    _bullets(s, 0.6, 2.0, 12, 4.5, research.get("insights", [])[:4])

    # 4 Market ---------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _section_header(s, "Market & Industry")
    _bar(s, 0.6, 1.4, 5.8, 2.2, LIGHT)
    _text(s, 0.8, 1.6, 5.4, 0.5, "Market size", size=14, color=MUTED)
    _text(s, 0.8, 2.1, 5.4, 1.2, f"${research.get('market_size_usd_bn','—')}B", size=44, bold=True, color=NAVY)
    _bar(s, 6.8, 1.4, 5.8, 2.2, LIGHT)
    _text(s, 7.0, 1.6, 5.4, 0.5, "CAGR (5Y)", size=14, color=MUTED)
    _text(s, 7.0, 2.1, 5.4, 1.2, f"{research.get('cagr_pct','—')}%", size=44, bold=True, color=GOLD)
    _text(s, 0.6, 4.0, 12, 0.5, "Key insights", size=16, bold=True, color=NAVY)
    _bullets(s, 0.6, 4.5, 12, 2.5, research.get("insights", [])[:4], size=14)

    # 5 Client snapshot ------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _section_header(s, f"Client Snapshot — {crm.get('client', client)}")
    rows = [
        ("Tier", crm.get("tier", "—")),
        ("RM Owner", crm.get("rm_owner", "—")),
        ("Wallet share", f"{crm.get('wallet_share_pct','—')}%"),
        ("Last meeting", crm.get("last_meeting", "—")),
        ("Open mandates", ", ".join(crm.get("open_mandates", []) or ["—"])),
        ("Products used", ", ".join(crm.get("products_used", []) or ["—"])),
    ]
    y = 1.4
    for label, value in rows:
        _text(s, 0.6, y, 3.2, 0.4, label, size=13, color=MUTED)
        _text(s, 3.9, y, 9, 0.4, str(value), size=15, bold=True, color=DARK)
        y += 0.55
    _text(s, 0.6, y + 0.2, 12, 0.4, "Key contacts", size=14, bold=True, color=NAVY)
    contacts = crm.get("key_contacts", [])
    _bullets(s, 0.6, y + 0.7, 12, 2, [f"{c.get('name')} — {c.get('role')}" for c in contacts])

    # 6 Competitors ----------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _section_header(s, "Competitor Landscape")
    peers = competitors.get("peers", [])
    col_w = 12 / max(len(peers), 1)
    for i, p in enumerate(peers):
        x = 0.6 + i * col_w
        _bar(s, x, 1.4, col_w - 0.2, 4.8, LIGHT)
        _bar(s, x, 1.4, col_w - 0.2, 0.5, NAVY)
        _text(s, x + 0.15, 1.45, col_w - 0.4, 0.5, p.get("name", ""), size=14, bold=True, color=RGBColor(255, 255, 255))
        _text(s, x + 0.15, 2.1, col_w - 0.4, 0.5, "Recent deal", size=11, color=MUTED)
        _text(s, x + 0.15, 2.5, col_w - 0.4, 1.5, p.get("recent_deal", ""), size=12, color=DARK)
        _text(s, x + 0.15, 4.4, col_w - 0.4, 0.5, "Strength", size=11, color=MUTED)
        _text(s, x + 0.15, 4.8, col_w - 0.4, 1.2, p.get("strength", ""), size=12, color=DARK)
    _text(s, 0.6, 6.4, 12, 0.5, f"Our edge: {competitors.get('our_edge','—')}", size=14, bold=True, color=GOLD)

    # 7 Financial profile ----------------------------------------------------
    s = prs.slides.add_slide(blank)
    _section_header(s, "Financial Profile & Comps")
    metrics = [
        ("Revenue ($M)", financials.get("revenue_usd_m", "—")),
        ("EBITDA margin", f"{financials.get('ebitda_margin_pct','—')}%"),
        ("Net leverage", f"{financials.get('net_leverage_x','—')}x"),
        ("EV / EBITDA", f"{financials.get('ev_ebitda_x','—')}x"),
    ]
    for i, (k, v) in enumerate(metrics):
        x = 0.6 + i * 3.1
        _bar(s, x, 1.4, 2.9, 1.6, LIGHT)
        _text(s, x + 0.2, 1.5, 2.6, 0.4, k, size=12, color=MUTED)
        _text(s, x + 0.2, 1.95, 2.6, 1.0, str(v), size=24, bold=True, color=NAVY)
    _text(s, 0.6, 3.4, 12, 0.5, "Selected precedent transactions", size=14, bold=True, color=NAVY)
    y = 3.95
    _text(s, 0.6, y, 5, 0.4, "Target", size=12, bold=True, color=MUTED)
    _text(s, 6.0, y, 3, 0.4, "EV ($M)", size=12, bold=True, color=MUTED)
    _text(s, 9.5, y, 3, 0.4, "EV / EBITDA", size=12, bold=True, color=MUTED)
    for d in financials.get("deal_comps", []):
        y += 0.5
        _text(s, 0.6, y, 5, 0.4, str(d.get("target", "—")), size=13)
        _text(s, 6.0, y, 3, 0.4, str(d.get("ev_usd_m", "—")), size=13)
        _text(s, 9.5, y, 3, 0.4, f"{d.get('ev_ebitda_x','—')}x", size=13)

    # 8 Strategic alternatives ----------------------------------------------
    s = prs.slides.add_slide(blank)
    _section_header(s, "Strategic Alternatives")
    alts = [
        "Status quo — optimize working capital and refinance 2027 maturities.",
        "Bolt-on M&A — 2-3 targets identified in adjacent segments.",
        "Strategic sale / partial monetization to a sponsor or strategic.",
        "IPO readiness pathway over 18-24 months.",
    ]
    _bullets(s, 0.6, 1.4, 12, 5.5, alts, size=18)

    # 9 Why us --------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _section_header(s, "Why Us")
    _bullets(
        s,
        0.6,
        1.4,
        12,
        5.5,
        [
            "Top-3 league table position in sector YTD.",
            f"{crm.get('wallet_share_pct','—')}% wallet share — long-standing trusted advisor.",
            "Integrated coverage: M&A, ECM, DCM, derivatives under one roof.",
            "Global distribution with deep sponsor relationships.",
        ],
        size=18,
    )

    # 10 Next steps ---------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _section_header(s, "Proposed Next Steps")
    _bullets(
        s,
        0.6,
        1.4,
        12,
        5.5,
        [
            "Working session with CFO & Treasurer to align on priorities.",
            "Deep-dive financial diagnostic (2 weeks).",
            "Refined strategic options memo & process timeline.",
            "Decision gate — mandate kickoff target end of quarter.",
        ],
        size=18,
    )

    # Closing --------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s, NAVY)
    _bar(s, 0.5, 3.5, 1.2, 0.08, GOLD)
    _text(s, 0.5, 3.7, 12, 0.8, "Thank you", size=44, bold=True, color=RGBColor(255, 255, 255))
    _text(s, 0.5, 4.6, 12, 0.5, "Confidential — for discussion purposes only", size=14, color=LIGHT)

    fname = f"pitchbook_{uuid.uuid4().hex[:10]}.pptx"
    out = os.path.join(OUT_DIR, fname)
    prs.save(out)
    return out
